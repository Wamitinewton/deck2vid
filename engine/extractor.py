from __future__ import annotations

from pathlib import Path
from typing import TypedDict

from pptx import Presentation
from pptx.shapes.base import BaseShape


class SlideContent(TypedDict):
    slide_number: int
    title: str
    body_text: list[str]
    notes: str
    # True when the slide contains no extractable text or speaker notes.
    # These slides are routed to the GPT-4.1 vision narration path.
    is_visual_only: bool


def _is_hidden(slide) -> bool:
    """Return True if the slide is marked hidden in the PPTX XML (show='0').

    LibreOffice silently excludes hidden slides when exporting to PDF, so we
    mirror that behaviour here to keep rendered image and audio counts in sync.
    """
    return slide._element.get("show") == "0"


def _extract_text_from_shape(shape: BaseShape) -> list[str]:
    texts: list[str] = []

    if shape.shape_type == 6:
        for s in shape.shapes:
            texts.extend(_extract_text_from_shape(s))
        return texts

    if not shape.has_text_frame:
        return texts

    for paragraph in shape.text_frame.paragraphs:
        line = paragraph.text.strip()
        if line:
            texts.append(line)

    return texts


def _extract_title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text.strip()
    return ""


def _extract_notes(slide) -> str:
    try:
        notes_frame = slide.notes_slide.notes_text_frame
        return notes_frame.text.strip()
    except Exception:
        return ""


def extract_slides(pptx_path: str | Path) -> list[SlideContent]:
    prs = Presentation(str(pptx_path))
    slides: list[SlideContent] = []

    # Enumerate with the raw PPTX slide number so slide_number values are
    # stable and meaningful in logs and output filenames.
    for index, slide in enumerate(prs.slides, start=1):
        if _is_hidden(slide):
            continue

        title = _extract_title(slide)
        body_text: list[str] = []

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            body_text.extend(_extract_text_from_shape(shape))

        notes = _extract_notes(slide)

        slides.append(
            SlideContent(
                slide_number=index,
                title=title,
                body_text=body_text,
                notes=notes,
                is_visual_only=not title and not body_text and not notes,
            )
        )

    return slides